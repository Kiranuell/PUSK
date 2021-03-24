from pptx import Presentation


prs = Presentation('text.pptx')

slides = [slide for slide in prs.slides]
shapes = []
itog = 0
sl = []
for slide in slides:
    shapes = []
    for shape in slide.shapes:
        shapes.append(shape)
    sl.append(shapes)

for i in sl:
    for j in i:
        if j.shape_type == 17:
            print(j.shape_type)
            print(j.font)
        elif j.shape_type == 1:
            print(j.shape_type)
            print(j.font)
        elif j.shape_type == 14:
            print(j.shape_type)
            print(j.font)
        else:
            print(j.shape_type)
    print("---")
print(itog)