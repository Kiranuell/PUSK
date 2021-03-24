from pptx import Presentation


prs = Presentation('background.pptx')

# print(len(prs.slides))
# print(prs.slides[0])
slides = [slide for slide in prs.slides]


for i in slides:
    print(i.background.fill.type)
    if i.background.fill.type == 1:
        print(i.background.fill.fore_color.type)
        if i.background.fill.fore_color.type == 1:
            print(i.background.fill.fore_color.rgb)
        if i.background.fill.fore_color.type == 2:
            print(i.background.fill.fore_color.brightness)
            print(i.background.fill.fore_color.theme_color)
    else:
        # print(i.background.fill.type)
        pass
    print("---")